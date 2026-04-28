import os
import logging
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import io
from openai import OpenAI
from dotenv import load_dotenv
from PIL import Image
import pytesseract
import json
from typing import Optional

load_dotenv()

logger = logging.getLogger("examai.engine")

MAX_QUESTION_COUNT = 100

_openrouter_api_key = os.environ.get("OPENROUTER_API_KEY")
if not _openrouter_api_key:
    raise EnvironmentError(
        "Missing required environment variable: OPENROUTER_API_KEY. "
        "Ensure it is set in your .env file or environment."
    )

# Setup OpenRouter client
client = OpenAI(
    api_key=os.environ.get("OPENROUTER_API_KEY"),
    base_url="https://openrouter.ai/api/v1",
    timeout=90.0,
)

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF, falling back to OCR for scanned pages."""
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text += page_text
            else:
                # Scanned page — render to image and OCR
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img, lang='heb+eng')
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX including OCR on embedded images."""
    doc = Document(io.BytesIO(file_bytes))
    parts = [paragraph.text for paragraph in doc.paragraphs]
    # OCR any inline images
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                img = Image.open(io.BytesIO(img_bytes))
                parts.append(pytesseract.image_to_string(img, lang='heb+eng'))
            except Exception as exc:
                logger.warning("DOCX image OCR failed for relation %s: %s", rel.reltype, exc)
    return "\n".join(parts)

def extract_text_from_pptx(file_bytes: bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

def extract_text_from_txt(file_bytes: bytes):
    return file_bytes.decode("utf-8")

def extract_text_from_image(file_bytes: bytes) -> str:
    """Extract text from JPG/PNG using OCR (Tesseract)."""
    image = Image.open(io.BytesIO(file_bytes))
    return pytesseract.image_to_string(image, lang='heb+eng')

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Route a single file to the correct extractor."""
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext == "pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == "txt":
        return extract_text_from_txt(file_bytes)
    elif ext in ("jpg", "jpeg", "png"):
        return extract_text_from_image(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _build_type_instruction(question_type: str, question_count: int) -> str:
    """Build the question-type section of the AI prompt."""
    type_instructions = {
        "open": f"""generate {question_count} open-ended questions in Hebrew.
Return a JSON object with a key "questions" containing an array of objects, each with:
- "question": the question text
- "answer": the correct answer in Hebrew
- "critical_points": an array of 3-5 key points in Hebrew that a complete answer must cover""",

        "yesno": f"""generate {question_count} yes/no questions in Hebrew. Each must be answerable with only כן or לא.
Return a JSON object with a key "questions" containing an array of objects, each with:
- "question": the question text
- "answer": either "כן" or "לא"

Important: distribute the correct answers randomly between כן and לא.
- Neither "כן" nor "לא" should be the correct answer for more than 70% of the questions.
- The correct answers must be spread randomly across all {question_count} questions.""",

        "multiple": f"""generate {question_count} multiple choice questions in Hebrew.
Return a JSON object with a key "questions" containing an array of objects, each with:
- "question": the question text only (no options)
- "options": an object with keys "א", "ב", "ג", "ד" each containing the option text in Hebrew
- "answer": the correct key, either "א", "ב", "ג", or "ד"

Important: distribute the correct answers randomly across all options.
- No single option ("א", "ב", "ג", or "ד") should be the correct answer for more than 50% of the questions.
- The correct answers must be spread randomly and naturally across all {question_count} questions.""",
    }
    return type_instructions[question_type]


def _build_difficulty_instruction(
    difficulty: str = "medium",
    distribution: Optional[dict] = None,
) -> str:
    """Build the difficulty section of the AI prompt.

    When *distribution* is supplied (a dict with 'easy', 'medium', 'hard' integer
    percentage keys summing to 100) it is used instead of the single *difficulty*
    string.
    """
    if distribution:
        easy = distribution.get("easy", 0)
        medium = distribution.get("medium", 0)
        hard = distribution.get("hard", 0)
        return (
            f"Distribute the difficulty of the questions as follows: "
            f"approximately {easy}% at the Remembering/Understanding level (Bloom's L1–L2), "
            f"{medium}% at the Applying/Analyzing level (Bloom's L3–L4), and "
            f"{hard}% at the Evaluating/Creating level (Bloom's L5–L6). "
            "Mix difficulty levels throughout the question set."
        )

    difficulty_instructions = {
        "easy": (
            "Focus specifically on the Remembering and Understanding levels of Bloom's Taxonomy (Levels 1–2). "
            "Questions should test the student's ability to recall facts, define terms, and explain basic concepts "
            "from the material. Avoid questions that require application or analysis."
        ),
        "medium": (
            "Focus specifically on the Applying and Analyzing levels of Bloom's Taxonomy (Levels 3–4). "
            "Questions should require the student to use information in new situations, solve problems using "
            "learned concepts, or draw connections and distinctions among ideas in the material."
        ),
        "hard": (
            "Focus specifically on the Evaluating and Creating levels of Bloom's Taxonomy (Levels 5–6). "
            "Questions should require the student to justify a position, critique a theory, synthesize information "
            "from multiple parts of the material, or propose an original solution or argument. "
            "Avoid simple recall or straightforward application questions."
        ),
    }
    return difficulty_instructions.get(difficulty, difficulty_instructions["medium"])


def _call_model(prompt: str) -> str:
    """Execute a single AI model call and return the raw JSON string."""
    response = client.chat.completions.create(
        model="openai/gpt-5.4-mini",
        messages=[
            {"role": "system", "content": "You are an expert educator. You must respond with valid JSON only."},
            {"role": "user", "content": prompt},
        ],
        stream=False,
        max_tokens=8192,
        timeout=90,
        response_format={"type": "json_object"}
    )
    return response.choices[0].message.content


def _build_prompt(
    text: str,
    type_instruction: str,
    difficulty_instruction: str,
    time_instruction: str = "",
) -> str:
    return f"""
    Analyze the following text and {type_instruction}
    {difficulty_instruction}
    Write all questions in Hebrew.
    {time_instruction}
    Before generating questions, check the content:
    - If the text is empty or too short to work with, return: {{"error": "הקובץ שהועלה ריק או קצר מדי. אנא העלה קובץ עם תוכן."}}
    - If the text contains random characters or meaningless content with no real information, return: {{"error": "התוכן בקובץ אינו משמעותי ולא ניתן ליצור ממנו שאלות. אנא העלה קובץ עם חומר לימוד תקין."}}
    - If the text is valid and meaningful, follow the format instructions above exactly.

    Text:
    {text}
    """


def _generate_mixed_questions(
    text: str,
    format_counts: dict,
    difficulty_instruction: str,
    include_time_estimation: bool,
) -> str:
    """Generate a mixed-type exam by running one batch per question type.

    Returns a JSON string with:
    - "questions": merged list with a "question_type" field on each item
    - "recommended_time" (int, minutes) if *include_time_estimation* is True
    """
    all_questions: list = []
    total_recommended_time = 0

    for q_type in ("yesno", "multiple", "open"):
        count = format_counts.get(q_type, 0)
        if count <= 0:
            continue

        time_instruction = ""
        if include_time_estimation:
            time_instruction = (
                f'Also add a "recommended_time" field (integer, in minutes) at the root level of the JSON, '
                f"estimating how long a student needs to complete these {count} {q_type} question(s) "
                "based on their content complexity and Bloom's levels."
            )

        prompt = _build_prompt(
            text,
            _build_type_instruction(q_type, count),
            difficulty_instruction,
            time_instruction,
        )
        raw = _call_model(prompt)
        try:
            batch = json.loads(raw)
        except json.JSONDecodeError:
            logger.warning("Failed to parse JSON for %s batch", q_type)
            continue

        if "error" in batch:
            return raw  # propagate the first error encountered

        for q in batch.get("questions", []):
            q["question_type"] = q_type
        all_questions.extend(batch.get("questions", []))

        if include_time_estimation:
            total_recommended_time += int(batch.get("recommended_time", 0))

    result: dict = {"questions": all_questions}
    if include_time_estimation:
        result["recommended_time"] = total_recommended_time
    return json.dumps(result, ensure_ascii=False)


def generate_questions(
    files: list[tuple[bytes, str]],
    question_type: str = "open",
    question_count: int = 5,
    difficulty: str = "medium",
    blueprint: Optional[dict] = None,
) -> str:
    """Generate exam questions from uploaded file content.

    Parameters
    ----------
    files:
        List of (file_bytes, filename) tuples.
    question_type:
        One of "open", "yesno", "multiple".  Ignored when *blueprint*
        contains non-trivial ``format_counts``.
    question_count:
        Total number of questions to generate (1–MAX_QUESTION_COUNT).
    difficulty:
        One of "easy", "medium", "hard".  Ignored when *blueprint*
        contains ``difficulty_distribution``.
    blueprint:
        Optional exam-blueprint configuration dict.  Supported keys:

        - ``time_mode`` – "manual" | "ai_estimated" (default "manual")
        - ``manual_time`` – integer minutes (informational only; stored in
          response when time_mode=="manual")
        - ``difficulty_distribution`` – dict with "easy", "medium", "hard"
          integer percentage values summing to 100
        - ``format_counts`` – dict with "yesno", "multiple", "open" integer
          counts.  When the values for more than one type are > 0 the exam
          is generated in mixed mode.
    """
    # 1. Extract and combine text from all uploaded files
    all_texts = []
    for file_bytes, filename in files:
        text = extract_text_from_file(file_bytes, filename)
        if text.strip():
            all_texts.append(f"[קובץ: {filename}]\n{text}")
    text = "\n\n---\n\n".join(all_texts)

    if question_count < 1:
        question_count = 1
    if question_count > MAX_QUESTION_COUNT:
        question_count = MAX_QUESTION_COUNT

    # ── Resolve blueprint options ──────────────────────────────────────────────
    bp = blueprint or {}
    time_mode = bp.get("time_mode", "manual")
    include_time_estimation = time_mode == "ai_estimated"
    distribution = bp.get("difficulty_distribution") or None
    format_counts = bp.get("format_counts") or {}

    # Count how many question types have a non-zero count requested
    active_types = [t for t in ("yesno", "multiple", "open") if format_counts.get(t, 0) > 0]
    is_mixed = len(active_types) > 1

    # ── Build difficulty instruction ───────────────────────────────────────────
    difficulty_instruction = _build_difficulty_instruction(difficulty, distribution)

    # ── Mixed exam: one batch per type ────────────────────────────────────────
    if is_mixed:
        return _generate_mixed_questions(text, format_counts, difficulty_instruction, include_time_estimation)

    # ── Single-type exam ──────────────────────────────────────────────────────
    # If format_counts specifies exactly one type, honour that count/type
    if len(active_types) == 1:
        question_type = active_types[0]
        question_count = format_counts[question_type]
        if question_count < 1:
            question_count = 1
        if question_count > MAX_QUESTION_COUNT:
            question_count = MAX_QUESTION_COUNT

    time_instruction = ""
    if include_time_estimation:
        time_instruction = (
            f'Also add a "recommended_time" field (integer, in minutes) at the root level of the JSON, '
            f"estimating how long a student needs to complete these {question_count} {question_type} question(s) "
            "based on their content complexity and Bloom's levels."
        )

    prompt = _build_prompt(
        text,
        _build_type_instruction(question_type, question_count),
        difficulty_instruction,
        time_instruction,
    )

    return _call_model(prompt)

def grade_answers(questions: list, answers: list, question_type: str):
    qa_text = ""
    for i, (q, a) in enumerate(zip(questions, answers)):
        # Use per-question type if available (mixed exams), fall back to the global type
        q_type = q.get("question_type", question_type)
        if q_type == "multiple":
            options_text = ", ".join([f"{k}: {v}" for k, v in q.get("options", {}).items()])
            qa_text += f"שאלה {i+1}: {q['question']}\nאפשרויות: {options_text}\nתשובה נכונה: {q['answer']}\nתשובת התלמיד: {a}\n\n"
        elif q_type == "yesno":
            qa_text += f"שאלה {i+1}: {q['question']}\nתשובה נכונה: {q['answer']}\nתשובת התלמיד: {a}\n\n"
        else:
            critical = "\n".join([f"- {p}" for p in q.get("critical_points", [])])
            qa_text += f"שאלה {i+1}: {q['question']}\nתשובה נכונה: {q['answer']}\nנקודות מפתח:\n{critical}\nתשובת התלמיד: {a}\n\n"

    prompt = f"""
    You are a strict expert educator grading a student's answers in Hebrew.
    Question type: {question_type}

    Here are the questions, correct answers, and the student's answers:
    {qa_text}

    Grading rules:
    - If the student's answer is random characters, gibberish, or meaningless text, it is ALWAYS wrong and gets 0 points.
    - For yes/no questions, only "כן" or "לא" are acceptable. Anything else gets 0 points.
    - For multiple choice questions, only the exact correct option key is acceptable. Anything else gets 0 points.
    - For open questions, grade based on the critical_points provided with each question:
        * If the student covers all critical points correctly: 1 point (full credit)
        * If the student covers some but not all critical points: 0.5 points (partial credit)
        * If the student does not cover any critical points or writes gibberish: 0 points
    - Be strict. Do not give credit for guesses or nonsense.

    Grade each answer and return a JSON object with:
    - "score": the total score as a number (can include 0.5 values for partial credit)
    - "feedback": an array of objects, one per question, each with:
        - "question": the question text
        - "points": the points awarded (0, 0.5, or 1)
        - "correct": true if points == 1, false otherwise
        - "covered_points": an array of the critical points the student covered correctly (in Hebrew)
        - "missed_points": an array of the critical points the student missed (in Hebrew)
        - "explanation": a short explanation in Hebrew of the grade

    Return ONLY valid JSON.
    """

    response = client.chat.completions.create(
        model="openai/gpt-5.4-mini",
        messages=[
            {"role": "system", "content": "You are an expert educator. You must respond with valid JSON only."},
            {"role": "user", "content": prompt},
        ],
        stream=False,
        max_tokens=8192,
        timeout=90,
        response_format={"type": "json_object"}
    )

    return response.choices[0].message.content