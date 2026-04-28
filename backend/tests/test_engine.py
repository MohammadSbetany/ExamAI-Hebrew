"""
Tests for engine.py — text extraction functions and input validation logic.
AI API calls are mocked so tests run without real credentials or network.
"""
import io
import json
import pytest
from unittest.mock import patch, MagicMock


# ── Mock OPENROUTER_API_KEY before importing engine ───────────────────────────
@pytest.fixture(autouse=True)
def mock_env(monkeypatch):
    monkeypatch.setenv("OPENROUTER_API_KEY", "test-key-mock")


# ── Text extraction ───────────────────────────────────────────────────────────

class TestExtractTextFromTxt:
    def test_basic_utf8(self):
        from engine import extract_text_from_txt
        content = "שלום עולם\nזהו טקסט לדוגמה".encode("utf-8")
        result = extract_text_from_txt(content)
        assert "שלום עולם" in result
        assert "טקסט לדוגמה" in result

    def test_empty_file(self):
        from engine import extract_text_from_txt
        result = extract_text_from_txt(b"")
        assert result == ""

    def test_english_text(self):
        from engine import extract_text_from_txt
        result = extract_text_from_txt(b"Hello World")
        assert result == "Hello World"


class TestExtractTextFromDocx:
    def test_valid_docx(self):
        from docx import Document
        from engine import extract_text_from_docx
        doc = Document()
        doc.add_paragraph("פסקה ראשונה")
        doc.add_paragraph("פסקה שנייה")
        buf = io.BytesIO()
        doc.save(buf)
        result = extract_text_from_docx(buf.getvalue())
        assert "פסקה ראשונה" in result
        assert "פסקה שנייה" in result

    def test_empty_docx(self):
        from docx import Document
        from engine import extract_text_from_docx
        doc = Document()
        buf = io.BytesIO()
        doc.save(buf)
        result = extract_text_from_docx(buf.getvalue())
        assert isinstance(result, str)


class TestExtractTextFromPptx:
    def test_valid_pptx(self):
        from pptx import Presentation
        from pptx.util import Inches
        from engine import extract_text_from_pptx
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "כותרת שקופית"
        slide.placeholders[1].text = "תוכן שקופית"
        buf = io.BytesIO()
        prs.save(buf)
        result = extract_text_from_pptx(buf.getvalue())
        assert "כותרת שקופית" in result
        assert "תוכן שקופית" in result

    def test_empty_pptx(self):
        from pptx import Presentation
        from engine import extract_text_from_pptx
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        result = extract_text_from_pptx(buf.getvalue())
        assert isinstance(result, str)


class TestExtractTextFromPdf:
    def test_valid_pdf(self):
        import fitz
        from engine import extract_text_from_pdf
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF content test")
        buf = io.BytesIO(doc.tobytes())
        result = extract_text_from_pdf(buf.getvalue())
        assert "PDF content test" in result

    def test_empty_pdf(self):
        import fitz
        from engine import extract_text_from_pdf
        doc = fitz.open()
        doc.new_page()
        buf = io.BytesIO(doc.tobytes())
        result = extract_text_from_pdf(buf.getvalue())
        assert isinstance(result, str)


# ── generate_questions ────────────────────────────────────────────────────────

class TestGenerateQuestions:
    def _mock_response(self, payload: dict):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_unsupported_file_type_raises(self):
        with patch.dict("os.environ", {"OPENROUTER_API_KEY": "test"}):
            from engine import generate_questions
            with pytest.raises(ValueError, match="Unsupported file type"):
                generate_questions([(b"content", "file.xyz")])

    def test_generates_open_questions(self):
        from engine import generate_questions
        expected = {"questions": [{"question": "מה זה?", "answer": "תשובה", "critical_points": ["נקודה"]}]}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(generate_questions([(b"some text", "test.txt")], "open", 1, "medium"))
        assert "questions" in result
        assert len(result["questions"]) == 1

    def test_generates_yesno_questions(self):
        from engine import generate_questions
        expected = {"questions": [{"question": "נכון?", "answer": "כן"}]}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(generate_questions([(b"some text", "test.txt")], "yesno", 1, "easy"))
        assert result["questions"][0]["answer"] in ("כן", "לא")

    def test_generates_multiple_choice_questions(self):
        from engine import generate_questions
        expected = {"questions": [{"question": "מה?", "answer": "א", "options": {"א": "א", "ב": "ב", "ג": "ג", "ד": "ד"}}]}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(generate_questions([(b"some text", "test.txt")], "multiple", 1, "hard"))
        assert "options" in result["questions"][0]

    def test_question_count_clamped_to_minimum(self):
        """question_count < 1 should be clamped to 1."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            # Should not raise even with count=0
            generate_questions([(b"text", "test.txt")], "open", 0, "medium")
            call_args = mock_client.chat.completions.create.call_args
            prompt = call_args[1]["messages"][1]["content"]
            assert "generate 1" in prompt

    def test_question_count_clamped_to_maximum(self):
        """question_count > MAX_QUESTION_COUNT should be clamped to MAX_QUESTION_COUNT."""
        from engine import generate_questions, MAX_QUESTION_COUNT
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"text", "test.txt")], "open", 999, "medium")
            call_args = mock_client.chat.completions.create.call_args
            prompt = call_args[1]["messages"][1]["content"]
            assert f"generate {MAX_QUESTION_COUNT}" in prompt

    def test_easy_difficulty_bloom_keywords_in_prompt(self):
        """Easy difficulty should reference Bloom's Taxonomy Levels 1-2 in the prompt."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"some text", "test.txt")], "open", 1, "easy")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "Bloom" in prompt
        assert "Levels 1" in prompt or "Level 1" in prompt

    def test_medium_difficulty_bloom_keywords_in_prompt(self):
        """Medium difficulty should reference Bloom's Taxonomy Levels 3-4 in the prompt."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"some text", "test.txt")], "open", 1, "medium")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "Bloom" in prompt
        assert "Levels 3" in prompt or "Level 3" in prompt

    def test_hard_difficulty_bloom_keywords_in_prompt(self):
        """Hard difficulty should reference Bloom's Taxonomy Levels 5-6 in the prompt."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"some text", "test.txt")], "open", 1, "hard")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "Bloom" in prompt
        assert "Levels 5" in prompt or "Level 5" in prompt

    def test_txt_file_dispatched_correctly(self):
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client, \
             patch("engine.extract_text_from_txt", return_value="extracted text") as mock_extract:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"hello", "test.txt")], "open", 1, "medium")
            mock_extract.assert_called_once_with(b"hello")

    def test_pdf_file_dispatched_correctly(self):
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client, \
             patch("engine.extract_text_from_pdf", return_value="pdf text") as mock_extract:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"pdf bytes", "doc.pdf")], "open", 1, "medium")
            mock_extract.assert_called_once()


# ── grade_answers ─────────────────────────────────────────────────────────────

class TestGradeAnswers:
    def _mock_response(self, payload: dict):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_grade_open_questions(self):
        from engine import grade_answers
        expected = {
            "score": 1,
            "feedback": [{"question": "מה?", "points": 1, "correct": True,
                          "explanation": "מצוין", "covered_points": ["נקודה"], "missed_points": []}]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(grade_answers(
                [{"question": "מה?", "answer": "תשובה", "critical_points": ["נקודה"]}],
                ["תשובה"],
                "open"
            ))
        assert result["score"] == 1

    def test_grade_yesno_questions(self):
        from engine import grade_answers
        expected = {
            "score": 1,
            "feedback": [{"question": "נכון?", "points": 1, "correct": True,
                          "explanation": "נכון", "covered_points": [], "missed_points": []}]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(grade_answers(
                [{"question": "נכון?", "answer": "כן"}],
                ["כן"],
                "yesno"
            ))
        assert result["score"] == 1

    def test_grade_multiple_choice(self):
        from engine import grade_answers
        expected = {
            "score": 0,
            "feedback": [{"question": "מה?", "points": 0, "correct": False,
                          "explanation": "שגוי", "covered_points": [], "missed_points": []}]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(grade_answers(
                [{"question": "מה?", "answer": "א", "options": {"א": "נ", "ב": "ל"}}],
                ["ב"],
                "multiple"
            ))
        assert result["score"] == 0

    def test_grade_partial_credit(self):
        from engine import grade_answers
        expected = {
            "score": 0.5,
            "feedback": [{"question": "שאלה?", "points": 0.5, "correct": False,
                          "explanation": "חלקי", "covered_points": ["נקודה 1"], "missed_points": ["נקודה 2"]}]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(grade_answers(
                [{"question": "שאלה?", "answer": "תשובה", "critical_points": ["נקודה 1", "נקודה 2"]}],
                ["תשובה חלקית"],
                "open"
            ))
        assert result["score"] == 0.5
        assert result["feedback"][0]["points"] == 0.5

class TestExtractTextFromImage:
    def test_image_ocr_called(self):
        """extract_text_from_image should call pytesseract."""
        from engine import extract_text_from_image
        from PIL import Image
        import io
        img = Image.new('RGB', (100, 30), color='white')
        buf = io.BytesIO()
        img.save(buf, format='PNG')
        with patch('engine.pytesseract.image_to_string', return_value='hello') as mock_ocr:
            result = extract_text_from_image(buf.getvalue())
            mock_ocr.assert_called_once()
            assert result == 'hello'

class TestExtractTextFromFile:
    def test_routes_jpg_to_image_extractor(self):
        from engine import extract_text_from_file
        with patch('engine.extract_text_from_image', return_value='img text') as mock:
            result = extract_text_from_file(b'bytes', 'photo.jpg')
            mock.assert_called_once_with(b'bytes')
            assert result == 'img text'

    def test_routes_png_to_image_extractor(self):
        from engine import extract_text_from_file
        with patch('engine.extract_text_from_image', return_value='img text') as mock:
            result = extract_text_from_file(b'bytes', 'photo.png')
            mock.assert_called_once_with(b'bytes')

    def test_unsupported_extension_raises(self):
        from engine import extract_text_from_file
        with pytest.raises(ValueError, match='Unsupported file type'):
            extract_text_from_file(b'bytes', 'file.xyz')


class TestSplitMergedQuestions:
    def test_split_sums_to_n(self):
        from engine import _split_merged_questions
        for n in range(1, 20):
            yn, mc, op = _split_merged_questions(n)
            assert yn + mc + op == n, f"Sum mismatch for n={n}: {yn}+{mc}+{op}"

    def test_split_open_at_least_one(self):
        from engine import _split_merged_questions
        for n in range(1, 10):
            _, _, op = _split_merged_questions(n)
            assert op >= 1, f"open_count < 1 for n={n}"

    def test_split_all_non_negative(self):
        from engine import _split_merged_questions
        for n in range(1, 20):
            yn, mc, op = _split_merged_questions(n)
            assert yn >= 0 and mc >= 0 and op >= 0


class TestMergedQuestionType:
    def _mock_response(self, payload):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_merged_type_instruction_references_all_types(self):
        """merged question_type prompt should mention yes/no, multiple, and open."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"text", "test.txt")], "merged", 10, "medium")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "yesno" in prompt or "yes/no" in prompt.lower()
        assert "multiple" in prompt
        assert "open" in prompt

    def test_merged_type_instruction_contains_type_field(self):
        """merged question_type prompt should instruct AI to include a 'type' field."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"text", "test.txt")], "merged", 10, "medium")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert '"type"' in prompt or "'type'" in prompt

    def test_merged_type_uses_split_counts_in_prompt(self):
        """merged prompt should include the computed split counts."""
        from engine import generate_questions, _split_merged_questions
        expected = {"questions": []}
        n = 10
        yn, mc, op = _split_merged_questions(n)
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"text", "test.txt")], "merged", n, "medium")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert str(yn) in prompt
        assert str(mc) in prompt
        assert str(op) in prompt


class TestMergedDifficulty:
    def _mock_response(self, payload):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_merged_difficulty_references_multiple_bloom_levels(self):
        """merged difficulty prompt should reference multiple Bloom levels."""
        from engine import generate_questions
        expected = {"questions": []}
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b"text", "test.txt")], "open", 5, "merged")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "Bloom" in prompt
        assert "L1" in prompt or "Levels 1" in prompt
        assert "L3" in prompt or "Levels 3" in prompt
        assert "L5" in prompt or "Levels 5" in prompt


class TestGradeAnswersMerged:
    def _mock_response(self, payload):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_grade_merged_uses_per_question_type_in_prompt(self):
        """grade_answers with merged type formats each question by its own type."""
        from engine import grade_answers
        questions = [
            {"question": "כן או לא?", "answer": "כן", "type": "yesno"},
            {"question": "מה?", "answer": "א",
             "options": {"א": "נ", "ב": "ל", "ג": "כ", "ד": "ד"}, "type": "multiple"},
            {"question": "הסבר", "answer": "תשובה", "critical_points": ["נקודה"], "type": "open"},
        ]
        answers = ["כן", "א", "תשובה"]
        expected = {
            "score": 3,
            "feedback": [
                {"question": "כן או לא?", "points": 1, "correct": True,
                 "explanation": "נכון", "covered_points": [], "missed_points": []},
                {"question": "מה?", "points": 1, "correct": True,
                 "explanation": "נכון", "covered_points": [], "missed_points": []},
                {"question": "הסבר", "points": 1, "correct": True,
                 "explanation": "נכון", "covered_points": ["נקודה"], "missed_points": []},
            ]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            result = json.loads(grade_answers(questions, answers, "merged"))
        assert result["score"] == 3
        prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        # Multiple choice question should have options in prompt
        assert "אפשרויות" in prompt

    def test_grade_merged_prompt_notes_mixed_type(self):
        """grade_answers with merged type should mention 'mixed' in the prompt."""
        from engine import grade_answers
        questions = [{"question": "שאלה?", "answer": "כן", "type": "yesno"}]
        expected = {
            "score": 1,
            "feedback": [{"question": "שאלה?", "points": 1, "correct": True,
                          "explanation": "נכון", "covered_points": [], "missed_points": []}]
        }
        with patch("engine.client") as mock_client:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            grade_answers(questions, ["כן"], "merged")
            prompt = mock_client.chat.completions.create.call_args[1]["messages"][1]["content"]
        assert "mixed" in prompt.lower() or "merged" in prompt.lower()

class TestMultipleFiles:
    def _mock_response(self, payload: dict):
        mock_choice = MagicMock()
        mock_choice.message.content = json.dumps(payload)
        mock_resp = MagicMock()
        mock_resp.choices = [mock_choice]
        return mock_resp

    def test_multiple_files_text_combined(self):
        """Text from all files should be combined with separator."""
        from engine import generate_questions
        expected = {'questions': []}
        with patch('engine.client') as mock_client, \
             patch('engine.extract_text_from_txt', return_value='text content') as mock_extract:
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b'a', 'a.txt'), (b'b', 'b.txt')], 'open', 1, 'medium')
            assert mock_extract.call_count == 2
            prompt = mock_client.chat.completions.create.call_args[1]['messages'][1]['content']
            assert '---' in prompt

    def test_empty_file_text_excluded(self):
        """Files that produce empty text should be excluded from prompt."""
        from engine import generate_questions
        expected = {'questions': []}
        with patch('engine.client') as mock_client, \
             patch('engine.extract_text_from_txt', side_effect=['', 'real content']):
            mock_client.chat.completions.create.return_value = self._mock_response(expected)
            generate_questions([(b'a', 'a.txt'), (b'b', 'b.txt')], 'open', 1, 'medium')
            prompt = mock_client.chat.completions.create.call_args[1]['messages'][1]['content']
            assert 'a.txt' not in prompt
            assert 'b.txt' in prompt    